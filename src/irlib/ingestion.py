"""File ingestion for common document formats.

The loader converts files and directories into `Document` records that can be
fed directly to any retriever.

Pseudocode:
    collect files from paths
    choose extractor by extension
    extract text plus source metadata
    optionally chunk long documents
"""

from __future__ import annotations

import csv
import json
import re
from pathlib import Path
from typing import Any, Iterable, Sequence

from irlib.core import Document, chunk_documents, read_text_lossy


class UnsupportedFormatError(ValueError):
    """Raised when a file cannot be parsed as a supported or text-like format."""


TEXT_EXTENSIONS = {
    ".txt",
    ".md",
    ".rst",
    ".py",
    ".js",
    ".ts",
    ".tsx",
    ".jsx",
    ".java",
    ".c",
    ".cpp",
    ".h",
    ".hpp",
    ".go",
    ".rs",
    ".rb",
    ".php",
    ".sql",
    ".yaml",
    ".yml",
    ".toml",
    ".ini",
    ".cfg",
}


def load_documents(
    path_or_paths: str | Path | Sequence[str | Path],
    *,
    recursive: bool = True,
    chunk: bool = True,
    chunk_size: int = 300,
    overlap: int = 50,
) -> list[Document]:
    """Load files or directories into `Document` objects.

    Pseudocode:
        paths = expand input files and directories
        docs = extract each file
        assign stable ids
        if chunk: split long docs into overlapping word windows

    Supported formats: txt, md, html, csv, tsv, json, jsonl, pdf, docx, xlsx,
    pptx, and common code/text files.
    """

    paths = list(_iter_paths(path_or_paths, recursive=recursive))
    raw_docs: list[Document] = []
    for path in paths:
        raw_docs.extend(_load_file(path))

    numbered = [
        Document(id=i, text=doc.text, metadata=doc.metadata, fields=doc.fields)
        for i, doc in enumerate(raw_docs)
        if doc.text.strip()
    ]
    if not chunk:
        return numbered
    return chunk_documents(numbered, chunk_size=chunk_size, overlap=overlap)


def _iter_paths(path_or_paths: str | Path | Sequence[str | Path], *, recursive: bool) -> Iterable[Path]:
    if isinstance(path_or_paths, (str, Path)):
        inputs = [Path(path_or_paths)]
    else:
        inputs = [Path(path) for path in path_or_paths]
    for input_path in inputs:
        if input_path.is_dir():
            iterator = input_path.rglob("*") if recursive else input_path.iterdir()
            for path in iterator:
                if path.is_file():
                    yield path
        elif input_path.is_file():
            yield input_path
        else:
            raise FileNotFoundError(str(input_path))


def _base_metadata(path: Path, **extra: Any) -> dict[str, Any]:
    metadata = {
        "source": str(path),
        "file_name": path.name,
        "extension": path.suffix.lower(),
    }
    metadata.update(extra)
    return metadata


def _doc(path: Path, text: str, **metadata: Any) -> Document:
    return Document(id=0, text=text, metadata=_base_metadata(path, **metadata))


def _load_file(path: Path) -> list[Document]:
    ext = path.suffix.lower()
    if ext in TEXT_EXTENSIONS:
        return [_doc(path, read_text_lossy(path))]
    if ext in {".html", ".htm"}:
        return [_doc(path, _load_html(path))]
    if ext in {".csv", ".tsv"}:
        return [_doc(path, _load_csv(path, delimiter="\t" if ext == ".tsv" else ","))]
    if ext == ".json":
        return _load_json(path)
    if ext == ".jsonl":
        return _load_jsonl(path)
    if ext == ".pdf":
        return _load_pdf(path)
    if ext == ".docx":
        return [_doc(path, _load_docx(path))]
    if ext == ".xlsx":
        return _load_xlsx(path)
    if ext == ".pptx":
        return _load_pptx(path)

    try:
        text = read_text_lossy(path)
    except Exception as exc:
        raise UnsupportedFormatError(f"Unsupported file format: {path}") from exc
    if not text.strip():
        raise UnsupportedFormatError(f"Unsupported file format: {path}")
    return [_doc(path, text)]


def _load_html(path: Path) -> str:
    html = read_text_lossy(path)
    try:
        from bs4 import BeautifulSoup

        return BeautifulSoup(html, "html.parser").get_text("\n")
    except Exception:
        return re.sub(r"<[^>]+>", " ", html)


def _load_csv(path: Path, *, delimiter: str) -> str:
    lines: list[str] = []
    with path.open("r", encoding="utf-8", newline="") as handle:
        reader = csv.reader(handle, delimiter=delimiter)
        for row in reader:
            lines.append(" ".join(cell for cell in row if cell is not None))
    return "\n".join(lines)


def _load_json(path: Path) -> list[Document]:
    value = json.loads(read_text_lossy(path))
    if isinstance(value, list):
        docs = []
        for i, item in enumerate(value):
            text = item.get("text", json.dumps(item, sort_keys=True)) if isinstance(item, dict) else json.dumps(item)
            docs.append(_doc(path, text, record=i))
        return docs
    text = value.get("text", json.dumps(value, sort_keys=True)) if isinstance(value, dict) else json.dumps(value)
    return [_doc(path, text)]


def _load_jsonl(path: Path) -> list[Document]:
    docs: list[Document] = []
    for i, line in enumerate(read_text_lossy(path).splitlines()):
        if not line.strip():
            continue
        value = json.loads(line)
        text = value.get("text", json.dumps(value, sort_keys=True)) if isinstance(value, dict) else json.dumps(value)
        docs.append(_doc(path, text, record=i))
    return docs


def _load_pdf(path: Path) -> list[Document]:
    try:
        from pypdf import PdfReader
    except Exception as exc:
        raise ImportError("Install `pypdf` to read PDF files.") from exc

    reader = PdfReader(str(path))
    docs = []
    for page_number, page in enumerate(reader.pages, start=1):
        text = page.extract_text() or ""
        docs.append(_doc(path, text, page=page_number))
    return docs


def _load_docx(path: Path) -> str:
    try:
        from docx import Document as DocxDocument
    except Exception as exc:
        raise ImportError("Install `python-docx` to read DOCX files.") from exc

    document = DocxDocument(str(path))
    parts = [paragraph.text for paragraph in document.paragraphs if paragraph.text]
    for table in document.tables:
        for row in table.rows:
            parts.append(" ".join(cell.text for cell in row.cells))
    return "\n".join(parts)


def _load_xlsx(path: Path) -> list[Document]:
    try:
        from openpyxl import load_workbook
    except Exception as exc:
        raise ImportError("Install `openpyxl` to read XLSX files.") from exc

    workbook = load_workbook(str(path), read_only=True, data_only=True)
    docs: list[Document] = []
    for sheet in workbook.worksheets:
        lines: list[str] = []
        for row in sheet.iter_rows(values_only=True):
            values = [str(value) for value in row if value is not None]
            if values:
                lines.append(" ".join(values))
        docs.append(_doc(path, "\n".join(lines), sheet=sheet.title))
    workbook.close()
    return docs


def _load_pptx(path: Path) -> list[Document]:
    try:
        from pptx import Presentation
    except Exception as exc:
        raise ImportError("Install `python-pptx` to read PPTX files.") from exc

    presentation = Presentation(str(path))
    docs: list[Document] = []
    for i, slide in enumerate(presentation.slides, start=1):
        parts: list[str] = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                parts.append(shape.text)
        docs.append(_doc(path, "\n".join(parts), slide=i))
    return docs


__all__ = ["UnsupportedFormatError", "load_documents"]

