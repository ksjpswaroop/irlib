import asyncio
import json

import pytest

from irlib import (
    AdvancedHybridRetriever,
    BM25FRetriever,
    BM25LRetriever,
    BM25PlusRetriever,
    BM25Retriever,
    BooleanRetriever,
    CachedRetriever,
    CharNGramRetriever,
    ColBERTInteraction,
    ConversationalQueryRewriter,
    DenseRetriever,
    DFRRetriever,
    ExactKNNRetriever,
    FacetedRetriever,
    FuzzyRetriever,
    HITSRanker,
    HashingEncoder,
    HNSWRetriever,
    HierarchicalRetriever,
    HybridRetriever,
    IVFRetriever,
    InvertedIndexRetriever,
    KnowledgeGraphRetriever,
    LIMEExplainer,
    LSHRetriever,
    LSIRetriever,
    LanguageModelRetriever,
    MMRDiversifier,
    PageRank,
    PairwiseLTRRanker,
    PermissionAwareRetriever,
    PersonalizedRetriever,
    PhraseRetriever,
    PointwiseLTRRanker,
    ProductQuantizationRetriever,
    ProximityRetriever,
    RAGRetriever,
    ReciprocalRankFusion,
    Reranker,
    SparseNeuralRetriever,
    StreamingDocumentProcessor,
    SynonymQueryExpander,
    TFIDFRetriever,
    TermFrequencyRetriever,
    TopicModelRetriever,
    WildcardPrefixRetriever,
    XQuADDiversifier,
    load_documents,
)


DOCS = [
    {"text": "Python retrieval search tutorial", "metadata": {"kind": "guide"}, "fields": {"title": "Python search", "body": "retrieval tutorial"}},
    {"text": "Java enterprise search platform", "metadata": {"kind": "guide"}, "fields": {"title": "Java search", "body": "enterprise platform"}},
    {"text": "Vector semantic retrieval for RAG", "metadata": {"kind": "paper"}, "fields": {"title": "Vector RAG", "body": "semantic retrieval"}},
]


def test_lexical_retrievers_return_ranked_results():
    classes = [
        InvertedIndexRetriever,
        BooleanRetriever,
        TermFrequencyRetriever,
        TFIDFRetriever,
        BM25Retriever,
        BM25PlusRetriever,
        BM25LRetriever,
        BM25FRetriever,
        LanguageModelRetriever,
        DFRRetriever,
        PhraseRetriever,
        ProximityRetriever,
        FuzzyRetriever,
        CharNGramRetriever,
        WildcardPrefixRetriever,
    ]
    for cls in classes:
        retriever = cls()
        retriever.index(DOCS)
        results = retriever.search("python retrieval", top_k=2)
        assert results, cls.__name__
        assert isinstance(results[0][0], int)
        assert isinstance(results[0][1], float)


def test_boolean_phrase_fuzzy_and_faceted_behavior():
    boolean = BooleanRetriever()
    boolean.index(DOCS)
    assert boolean.search("python AND retrieval")[0][0] == 0
    assert all(doc_id != 1 for doc_id, _ in boolean.search("python NOT java"))

    phrase = PhraseRetriever()
    phrase.index(DOCS)
    assert phrase.search('"python retrieval"')[0][0] == 0

    fuzzy = FuzzyRetriever(threshold=0.75)
    fuzzy.index(DOCS)
    assert fuzzy.search("pythn")[0][0] == 0

    faceted = FacetedRetriever(filters={"kind": "paper"})
    faceted.index(DOCS)
    assert faceted.search("retrieval")[0][0] == 2


def test_dense_ann_and_semantic_retrievers_use_fallback_encoder():
    retrievers = [
        DenseRetriever(encoder=HashingEncoder()),
        ExactKNNRetriever(encoder=HashingEncoder()),
        LSHRetriever(encoder=HashingEncoder()),
        HNSWRetriever(encoder=HashingEncoder()),
        IVFRetriever(encoder=HashingEncoder()),
        ProductQuantizationRetriever(encoder=HashingEncoder()),
        LSIRetriever(),
        TopicModelRetriever(),
        SparseNeuralRetriever(),
    ]
    for retriever in retrievers:
        retriever.index(DOCS)
        assert retriever.search("semantic retrieval", top_k=2), type(retriever).__name__


def test_hybrid_rerank_and_advanced_helpers():
    hybrid = HybridRetriever(dense_retriever=DenseRetriever(encoder=HashingEncoder()))
    hybrid.index(DOCS)
    candidates = hybrid.search("semantic retrieval", top_k=3)
    assert candidates

    reranker = Reranker(scorer=lambda q, p: 1.0 if "semantic" in p.lower() else 0.0)
    passages = [hybrid.get_document(doc_id).text for doc_id, _ in candidates]
    assert reranker.rerank("semantic retrieval", passages, [doc_id for doc_id, _ in candidates])[0][0] == 2

    advanced = AdvancedHybridRetriever(reranker=reranker, dense_retriever=DenseRetriever(encoder=HashingEncoder()))
    advanced.index(DOCS)
    assert advanced.search("semantic retrieval", top_k=1)[0][0] == 2

    colbert = ColBERTInteraction(encoder=HashingEncoder())
    assert colbert.compute_late_interaction("semantic retrieval", "semantic retrieval for rag") > 0

    explainer = LIMEExplainer(lambda q, d: float(len(set(q.split()) & set(d.split()))))
    explanation = explainer.explain("semantic retrieval", "semantic retrieval")
    assert explanation["impact"]


def test_ranking_expansion_diversification_and_graph_utilities():
    expander = SynonymQueryExpander({"python": ["programming"]})
    assert "programming" in expander.expand("python search")

    rrf = ReciprocalRankFusion()
    assert rrf.fuse([[(1, 2.0), (2, 1.0)], [(2, 3.0)]], top_k=1)[0][0] == 2

    mmr = MMRDiversifier(lambda_mult=0.5)
    selected = mmr.select([1, 2], {1: 1.0, 2: 0.9}, similarity=lambda a, b: 1.0 if a == b else 0.0, top_k=2)
    assert selected == [1, 2]

    xquad = XQuADDiversifier()
    assert xquad.select([1, 2], {1: 0.7, 2: 0.6}, {"a": {2: 1.0}}, top_k=1)

    graph = {"A": ["B", "C"], "B": ["C"], "C": ["A"]}
    assert PageRank().rank(graph)
    authorities, hubs = HITSRanker().rank(graph)
    assert authorities and hubs
    assert KnowledgeGraphRetriever(graph).search("A")

    rewriter = ConversationalQueryRewriter()
    assert "previous topic" in rewriter.rewrite(["previous topic"], "what about latency")


def test_ltr_permission_personalized_rag_and_hierarchical():
    pointwise = PointwiseLTRRanker(epochs=5).fit([[1.0, 0.0], [0.0, 1.0]], [1.0, 0.0])
    assert pointwise.rank([[1.0, 0.0], [0.0, 1.0]], top_k=1)[0][0] == 0

    pairwise = PairwiseLTRRanker(epochs=3).fit([([1.0, 0.0], [0.0, 1.0])])
    assert pairwise.rank([[1.0, 0.0], [0.0, 1.0]], top_k=1)[0][0] == 0

    base = BM25Retriever()
    permitted = PermissionAwareRetriever(base, lambda user, doc: doc.metadata["kind"] == user)
    permitted.index(DOCS)
    assert permitted.search("retrieval", user="paper")[0][0] == 2

    personalized = PersonalizedRetriever(base, lambda context, doc: 10.0 if doc.metadata["kind"] == context.get("kind") else 0.0)
    personalized.index(DOCS)
    assert personalized.search("retrieval", context={"kind": "paper"})[0][0] == 2

    rag = RAGRetriever(HybridRetriever(dense_retriever=DenseRetriever(encoder=HashingEncoder())))
    rag.index(DOCS)
    assert "semantic" in rag.retrieve_context("semantic retrieval", top_k=1).lower()

    hierarchical = HierarchicalRetriever(chunk_size=3, overlap=1)
    hierarchical.index(DOCS)
    assert hierarchical.search("semantic retrieval")


def test_async_cache_and_streaming():
    retriever = BM25Retriever()
    docs = StreamingDocumentProcessor(batch_size=2).process([doc["text"] for doc in DOCS], retriever=retriever)
    assert len(docs) == 3

    cached = CachedRetriever(retriever)
    first = asyncio.run(cached.search("python", top_k=1))
    second = asyncio.run(cached.search("python", top_k=1))
    assert first == second


def test_load_documents_text_json_csv_html(tmp_path):
    (tmp_path / "note.md").write_text("Python markdown retrieval", encoding="utf-8")
    (tmp_path / "page.html").write_text("<html><body>HTML search content</body></html>", encoding="utf-8")
    (tmp_path / "rows.csv").write_text("name,body\nalpha,csv retrieval\n", encoding="utf-8")
    (tmp_path / "records.jsonl").write_text(json.dumps({"text": "jsonl semantic retrieval"}) + "\n", encoding="utf-8")

    docs = load_documents(tmp_path, chunk=False)
    text = "\n".join(doc.text for doc in docs)
    assert "markdown retrieval" in text
    assert "HTML search content" in text
    assert "csv retrieval" in text
    assert "jsonl semantic retrieval" in text
    assert all("source" in doc.metadata for doc in docs)


def test_load_documents_generated_office_and_pdf_files(tmp_path):
    reportlab = pytest.importorskip("reportlab.pdfgen.canvas")
    pytest.importorskip("docx")
    pytest.importorskip("openpyxl")
    pytest.importorskip("pptx")

    from docx import Document as DocxDocument
    from openpyxl import Workbook
    from pptx import Presentation
    from reportlab.pdfgen import canvas

    pdf_path = tmp_path / "sample.pdf"
    pdf = canvas.Canvas(str(pdf_path))
    pdf.drawString(100, 750, "PDF retrieval content")
    pdf.save()

    docx_path = tmp_path / "sample.docx"
    docx = DocxDocument()
    docx.add_paragraph("DOCX retrieval content")
    docx.save(docx_path)

    xlsx_path = tmp_path / "sample.xlsx"
    workbook = Workbook()
    sheet = workbook.active
    sheet.title = "Data"
    sheet.append(["name", "body"])
    sheet.append(["alpha", "XLSX retrieval content"])
    workbook.save(xlsx_path)

    pptx_path = tmp_path / "sample.pptx"
    deck = Presentation()
    slide = deck.slides.add_slide(deck.slide_layouts[5])
    slide.shapes.title.text = "PPTX retrieval content"
    deck.save(pptx_path)

    docs = load_documents([pdf_path, docx_path, xlsx_path, pptx_path], chunk=False)
    text = "\n".join(doc.text for doc in docs)
    assert "PDF retrieval content" in text
    assert "DOCX retrieval content" in text
    assert "XLSX retrieval content" in text
    assert "PPTX retrieval content" in text
    assert any(doc.metadata.get("page") == 1 for doc in docs)
    assert any(doc.metadata.get("sheet") == "Data" for doc in docs)
    assert any(doc.metadata.get("slide") == 1 for doc in docs)
